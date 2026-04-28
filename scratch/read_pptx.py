import collections
import collections.abc
collections.Iterator = collections.abc.Iterator
collections.Iterable = collections.abc.Iterable
collections.Mapping = collections.abc.Mapping
collections.MutableSet = collections.abc.MutableSet
collections.MutableMapping = collections.abc.MutableMapping
collections.MutableSequence = collections.abc.MutableSequence

from pptx import Presentation

prs = Presentation('GIANT_APRESENTACAO_PREMIUM.pptx')
for i, slide in enumerate(prs.slides):
    if i >= 3: break
    print(f"--- Slide {i} ---")
    for shape in slide.shapes:
        text = shape.text if hasattr(shape, "text") else ""
        print(f"Shape: {shape.shape_type}, Name: {shape.name}, Text: {repr(text)}")
