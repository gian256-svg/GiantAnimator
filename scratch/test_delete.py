import collections
import collections.abc
collections.Iterator = collections.abc.Iterator
collections.Iterable = collections.abc.Iterable
collections.Mapping = collections.abc.Mapping
collections.MutableSet = collections.abc.MutableSet
collections.MutableMapping = collections.abc.MutableMapping
collections.MutableSequence = collections.abc.MutableSequence

from pptx import Presentation

def delete_slide(prs, slide):
    # Find the slide id
    id_dict = { slide.id: [i, slide.rId] for i,slide in enumerate(prs.slides._sldIdLst) }
    slide_id = slide.slide_id
    prs.part.drop_rel(id_dict[slide_id][1])
    del prs.slides._sldIdLst[id_dict[slide_id][0]]

prs = Presentation('GIANT_APRESENTACAO_PREMIUM.pptx')
print(f"Original slides: {len(prs.slides)}")
while len(prs.slides) > 2:
    delete_slide(prs, prs.slides[-1])
print(f"New slides: {len(prs.slides)}")
prs.save('scratch/test_delete.pptx')
