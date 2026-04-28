import collections
import collections.abc
collections.Iterator = collections.abc.Iterator
collections.Iterable = collections.abc.Iterable
collections.Mapping = collections.abc.Mapping
collections.MutableSet = collections.abc.MutableSet
collections.MutableMapping = collections.abc.MutableMapping
collections.MutableSequence = collections.abc.MutableSequence

from pptx import Presentation
from pptx.dml.color import RGBColor

SLIDES_CONTENT = [
    # Slide 0: Capa
    None,
    
    # Slide 1: Introdução
    {
        "title": "GiantAnimator: Do Estático ao Autônomo",
        "body": "Apresentação Executiva - Abril 2026\n\nNesta apresentação, exploraremos a arquitetura, os diferenciais e a oportunidade de mercado do GiantAnimator.\nVamos mergulhar no motor de renderização 4K mais preciso do mercado financeiro."
    },
    
    # Slide 2: O Desafio & Solução
    {
        "title": "O Desafio da Escala vs. Solução Nativa 4K",
        "body": "• O volume massivo de dados financeiros em 2026 exige respostas instantâneas.\n• Relatórios estáticos tradicionais estão obsoletos e não engajam investidores modernos.\n• Existe um vácuo imenso entre a análise de dados complexa e a comunicação visual.\n• A Solução: Nosso Motor de vídeo nativo UHD foi construído para os desafios da próxima década.\n• Resolvemos a latência na geração de vídeos com renderização em tempo real via React/Remotion."
    },
    
    # Slide 3: Diferencial Técnico
    {
        "title": "Fluxo Cirúrgico & Silent Auditor Loop",
        "body": "• O diferencial absoluto do GiantAnimator é seu \"Surgery-Grade Flow\" operado por IA de Visão.\n• Feedback Loop Autônomo garante 100% de fidelidade nos dados renderizados.\n• O sistema realiza a auto-correção de eixos e labels em tempo real, sem intervenção humana.\n• Pipeline resiliente capaz de detectar inconsistências numéricas (clipping, overlap) e corrigi-las.\n• Garantia matemática de que a barra exibida corresponde perfeitamente ao valor numérico exigido."
    },
    
    # Slide 4: Eficiência
    {
        "title": "Eficiência de Custos em Larga Escala",
        "body": "• Infraestrutura 100% escalável que garante a materialização de conteúdos UHD dinâmicos.\n• Análise cirúrgica otimizada a um custo marginal < US$ 0,001 por análise gráfica.\n• Arquitetura baseada no Gemini Flash Pricing, viabilizando operações em lote (Batch Processing).\n• Eliminação da necessidade de designers em tarefas repetitivas de diagramação financeira.\n• Redução do tempo de entrega (SLA) de semanas para meros segundos."
    },
    
    # Slide 5: Ecossistema
    {
        "title": "Integração Nativa & Biblioteca de Ativos",
        "body": "• Injeção direta e transparente de planilhas Excel e arquivos CSV.\n• O pipeline automatizado mitiga agressivamente falhas de erro humano (o temido \"fat-finger\").\n• Biblioteca de Ativos Treinados já inclui mais de 20 modelos de gráficos corporativos.\n• Suporte nativo a Waterfall, Heatmaps, Treemaps, Bar, Line, Radar e Gauge Charts.\n• Expansão contínua com estilos visuais hiper-personalizáveis (Theme Intelligence)."
    },
    
    # Slide 6: Mercado
    {
        "title": "Potencial de Mercado: O 'Blue Ocean'",
        "body": "• Setor B2B (Investor Relations & Wealth): Vídeos trimestrais automatizados para reporte de lucros.\n• Relatórios VIP hiper-personalizados sob demanda para clientes de alta renda (High-Net-Worth).\n• SaaS Pro: Assinatura escalável focada em analistas financeiros e criadores de conteúdo de dados.\n• Enterprise Solution: Instalação On-Premise 100% isolada e segura, ideal para Bancos e Corretoras.\n• Um mercado inexplorado que une Data Science puro com Motion Graphics de altíssimo nível."
    },
    
    # Slide 7: Roadmap
    {
        "title": "Visão de Futuro Estratégica (2026 - 2027)",
        "body": "• Q3/2026: Lançamento oficial da integração de Narração Neural (TTS) hiper-realista.\n• Sincronização automática de voz com animações e destaques visuais.\n• Q1/2027: Lançamento dos revolucionários \"Multi-Agent Reviewers\".\n• Implementação de Consenso de Dados Matemáticos Descentralizado para validação de métricas.\n• 2030: Domínio do mercado global de automação de vídeos institucionais baseados em dados."
    },
    
    # Slide 8: TI Segurança
    {
        "title": "Governança de Dados & Privacidade (Zero-Trust)",
        "body": "• Operação On-Premise / Local: Os datasets confidenciais são processados inteiramente em sandbox local.\n• Interação Headless CLI: Arquitetura server-side sem UI web exposta externamente.\n• Mitigação profunda contra riscos de Injeção de Código (SQLi/XSS) e Ataques de DDoS.\n• As chamadas de IA transmitem exclusivamente metadados ou frames ofuscados (pixels sem dados sensíveis).\n• Conformidade rigorosa: Zero tráfego de PII (Personally Identifiable Information) para LLMs externos."
    },
    
    # Slide 9: Mitigação Riscos
    {
        "title": "Mitigação de Riscos & Alucinação (AI Safety)",
        "body": "• O grande risco corporativo: A maioria das IAs gerativas de vídeo sofre de \"Data Hallucination\".\n• Como resolvemos: O vídeo obedece a cálculos nativos determinísticos (Remotion/React).\n• A IA do GiantAnimator atua EXCLUSIVAMENTE como auditora (Fiscal de Pixels), não como geradora de frames.\n• Resiliência de Infraestrutura: Sistema construído com Fallback Loop avançado.\n• Em caso de queda da API de IA, o sistema garante a renderização via algoritmos fallback 100% exatos."
    },
    
    # Slide 10: Encerramento
    {
        "title": "GiantAnimator",
        "body": "A Revolução Autônoma de Dados Visuais.\n\nObrigado.\n2026 © GiantAnimator Systems."
    }
]

def update_text_frame(shape, new_text, is_body=False):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    
    # 1. ATIVAR WORD WRAP para não sair da margem!
    tf.word_wrap = True
    
    if len(tf.paragraphs) == 0:
        return
    
    # Salva formatação da primeira run
    p0 = tf.paragraphs[0]
    font_name = None
    font_size = None
    font_color = None
    font_bold = None
    if len(p0.runs) > 0:
        r0 = p0.runs[0]
        font_name = r0.font.name
        font_size = r0.font.size
        try:
            font_color = r0.font.color.rgb if hasattr(r0.font.color, 'rgb') else None
        except:
            font_color = None
        font_bold = r0.font.bold

    tf.clear()
    
    # Quebra as linhas e adiciona parágrafos mantendo formatação
    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        
        # Aplica a fonte salva em cada run do parágrafo
        for r in p.runs:
            if font_name: r.font.name = font_name
            if font_size: r.font.size = font_size
            if font_bold is not None: r.font.bold = font_bold
            
            if is_body:
                r.font.color.rgb = RGBColor(255, 255, 255)
            elif font_color:
                try:
                    r.font.color.rgb = font_color
                except:
                    pass

def main():
    # Carregamos o arquivo original (não o corrompido) ou o atual se já estiver limpo.
    # Mas como o test_delete.py e o script anterior alteraram o arquivo, vamos tentar abrir.
    # Se o arquivo tiver menos de 11 slides, a gente tenta avisar, mas o user tinha o arquivo completo antes.
    # Como eu editei na mesma pasta, o arquivo já tem 11 slides.
    prs = Presentation('GIANT_APRESENTACAO_PREMIUM.pptx')
    
    for i in range(len(SLIDES_CONTENT)):
        if i >= len(prs.slides):
            break
            
        content = SLIDES_CONTENT[i]
        if content is None:
            continue
            
        slide = prs.slides[i]
        
        # Encontra text boxes
        text_shapes = [shape for shape in slide.shapes if shape.has_text_frame]
        
        if len(text_shapes) >= 1:
            update_text_frame(text_shapes[0], content["title"], is_body=False)
        if len(text_shapes) >= 2:
            update_text_frame(text_shapes[1], content["body"], is_body=True)

    prs.save('GIANT_APRESENTACAO_PREMIUM.pptx')
    print("Apresentação expandida e margens corrigidas com sucesso.")

if __name__ == '__main__':
    main()
